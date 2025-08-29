from typing import Dict, Any, List, Optional
import os
import re
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE


def _load_template() -> "Presentation":
    if os.path.exists("template.pptx"):
        try:
            return Presentation("template.pptx")
        except Exception:
            pass
    return Presentation()


def _first_body_placeholder(slide) -> Optional[object]:
    for shp in slide.placeholders:
        try:
            if shp.placeholder_format.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.CONTENT, PP_PLACEHOLDER.OBJECT):
                return shp
        except Exception:
            continue
    return None


def _add_textbox(slide, left=Inches(0.5), top=Inches(1.5), width=Inches(9), height=Inches(4.5)):
    return slide.shapes.add_textbox(left, top, width, height)


def _clear_text_frame(tf):
    try:
        for p in list(tf.paragraphs):
            tf._element.remove(p._p)
    except Exception:
        pass


def _split_overview(text: str, max_lines: int = 3) -> List[str]:
    sentences = re.split(r"[.?!]", text)
    sentences = [s.strip() for s in sentences if s.strip()]
    return sentences[:max_lines] if sentences else [text]


def _render_chart(slide, chart_spec, left=Inches(5.5), top=Inches(2), width=Inches(4), height=Inches(3)):
    try:
        categories = chart_spec.get("categories", [])
        values = chart_spec.get("values", [])

        # Auto-fix: generate categories if missing
        if values and not categories:
            categories = [f"Item {i+1}" for i in range(len(values))]

        # Only render if valid
        if not categories or not values or len(categories) != len(values):
            print("[WARN] Skipping chart: invalid categories/values")
            return

        data = CategoryChartData()
        data.categories = categories
        data.add_series(chart_spec.get("title", "Series"), values)

        chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED if chart_spec.get("type") == "bar" else XL_CHART_TYPE.LINE
        slide.shapes.add_chart(chart_type, left, top, width, height, data)
    except Exception as e:
        print("[WARN] Failed to render chart:", e)



def _render_table(slide, table_spec, left=Inches(0.5), top=Inches(4.5), width=Inches(9), height=Inches(2)):
    try:
        rows = len(table_spec.get("rows", [])) + 1
        cols = len(table_spec.get("headers", []))
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table

        for j, h in enumerate(table_spec.get("headers", [])):
            table.cell(0, j).text = str(h)

        for i, row in enumerate(table_spec.get("rows", []), start=1):
            for j, val in enumerate(row):
                table.cell(i, j).text = str(val)
    except Exception as e:
        print("[WARN] Failed to render table:", e)


def _render_topics(slide, body_shape, topics: List[Dict[str, Any]]):
    if body_shape is None or not hasattr(body_shape, "text_frame"):
        return
    tf = body_shape.text_frame
    _clear_text_frame(tf)

    for i, topic in enumerate(topics):
        subtitle = topic.get("subtitle", "")
        bullets = topic.get("bullets", [])
        sources = topic.get("sources", [])

        # Subtitle
        if subtitle:
            p = tf.add_paragraph()
            p.text = subtitle
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(20)
                r.font.color.rgb = RGBColor(0, 51, 102)

        # Bullets
        for b in bullets:
            p = tf.add_paragraph()
            p.text = str(b)
            p.level = 1
            for r in p.runs:
                r.font.size = Pt(16)
                r.font.color.rgb = RGBColor(40, 40, 40)

        # Sources
        if sources:
            src_text = "Sources: " + "; ".join(sources)
            p = tf.add_paragraph()
            p.text = src_text
            for r in p.runs:
                r.font.size = Pt(10)
                r.font.italic = True
                r.font.color.rgb = RGBColor(120, 120, 120)

        # Chart
        if "chart" in topic:
            _render_chart(slide, topic["chart"])

        # Table
        if "table" in topic:
            _render_table(slide, topic["table"])

        if i < len(topics) - 1:
            tf.add_paragraph().text = ""


def _safe_set_title(slide, text: str):
    try:
        if getattr(slide.shapes, "title", None):
            slide.shapes.title.text = text
            for p in slide.shapes.title.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(28)
                    r.font.color.rgb = RGBColor(0, 51, 102)
        else:
            tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(1))
            tf = tx.text_frame
            _clear_text_frame(tf)
            p = tf.add_paragraph()
            p.text = text
            for r in p.runs:
                r.font.size = Pt(28)
    except Exception:
        pass


def create_presentation(content: Dict[str, Any], output_filename: str):
    prs = _load_template()
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]

    # Title Slide
    s0 = prs.slides.add_slide(title_layout)
    _safe_set_title(s0, content.get("title", "Untitled Presentation"))

    # Overview
    s1 = prs.slides.add_slide(content_layout)
    _safe_set_title(s1, "Overview")
    body = _first_body_placeholder(s1) or _add_textbox(s1)
    overview_lines = _split_overview(content.get("overview", ""), max_lines=3)
    _render_topics(s1, body, [{"subtitle": "", "bullets": overview_lines, "sources": []}])

    # Content Slides
    for slide_data in content.get("slides", []):
        slide = prs.slides.add_slide(content_layout)
        _safe_set_title(slide, slide_data.get("title", "Slide"))
        body = _first_body_placeholder(slide) or _add_textbox(slide)
        topics = slide_data.get("topics", [])
        if not isinstance(topics, list):
            topics = [topics]
        _render_topics(slide, body, topics)

    # Conclusion
    if content.get("conclusion"):
        sl = prs.slides.add_slide(content_layout)
        _safe_set_title(sl, "Conclusion")
        body = _first_body_placeholder(sl) or _add_textbox(sl)
        _render_topics(sl, body, [{"subtitle": "", "bullets": [content["conclusion"]], "sources": []}])

    if not output_filename.lower().endswith(".pptx"):
        output_filename += ".pptx"
    prs.save(output_filename)
